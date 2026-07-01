from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image
import os

HERE = os.path.dirname(os.path.abspath(__file__))
OUT = os.path.join(os.path.dirname(HERE), "Demo-Payroll-System-Pitch.pptx")

NAVY = RGBColor(0x0F, 0x1B, 0x33)
NAVY_LIGHT = RGBColor(0x1B, 0x2A, 0x4A)
AMBER = RGBColor(0xF2, 0xA4, 0x1D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x5B, 0x66, 0x77)
LIGHT_BG = RGBColor(0xF6, 0xF7, 0xF9)
GREEN = RGBColor(0x1E, 0x8E, 0x5A)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def add_slide():
    return prs.slides.add_slide(BLANK)


def fill_bg(slide, color):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    spTree = slide.shapes._spTree
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)
    return bg


def add_text(slide, left, top, width, height, text, size=18, color=NAVY, bold=False,
             align=PP_ALIGN.LEFT, font="Calibri", line_spacing=1.0, anchor=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    if anchor:
        tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        p.line_spacing = line_spacing
        for r in p.runs:
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color
            r.font.name = font
    return box


def add_bullets(slide, left, top, width, height, items, size=15, color=NAVY,
                 bullet_color=AMBER, gap=6, bold_lead=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        run_bullet = p.add_run()
        run_bullet.text = "■  "
        run_bullet.font.size = Pt(size)
        run_bullet.font.color.rgb = bullet_color
        run_bullet.font.bold = True
        if isinstance(item, tuple):
            lead, rest = item
            r1 = p.add_run()
            r1.text = lead
            r1.font.size = Pt(size)
            r1.font.bold = True
            r1.font.color.rgb = color
            r2 = p.add_run()
            r2.text = rest
            r2.font.size = Pt(size)
            r2.font.color.rgb = color
        else:
            r = p.add_run()
            r.text = item
            r.font.size = Pt(size)
            r.font.color.rgb = color
    return box


def add_kicker(slide, text, color=AMBER, left=Inches(0.6), top=Inches(0.35)):
    add_text(slide, left, top, Inches(6), Inches(0.4), text.upper(), size=13, color=color, bold=True)


def add_pagenum(slide, n):
    add_text(slide, SW - Inches(1.0), SH - Inches(0.5), Inches(0.6), Inches(0.35),
              str(n), size=11, color=GRAY, align=PP_ALIGN.RIGHT)


def picture_framed(slide, path, left, top, max_w, max_h, border=True):
    im = Image.open(path)
    iw, ih = im.size
    ratio = min(max_w / iw, max_h / ih)
    w, h = int(iw * ratio), int(ih * ratio)
    l = int(left + (max_w - w) / 2)
    t = int(top + (max_h - h) / 2)
    if border:
        pad = Emu(int(Inches(0.05)))
        frame = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l - pad, t - pad, w + 2 * pad, h + 2 * pad)
        frame.fill.solid()
        frame.fill.fore_color.rgb = WHITE
        frame.line.color.rgb = RGBColor(0xD8, 0xDC, 0xE2)
        frame.line.width = Pt(0.75)
        frame.shadow.inherit = False
    slide.shapes.add_picture(path, l, t, width=w, height=h)


IMG = lambda name: os.path.join(HERE, name)

# ---------------------------------------------------------------- Slide 1: Title
s = add_slide()
fill_bg(s, NAVY)
accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(6.55), SW, Inches(0.12))
accent.fill.solid(); accent.fill.fore_color.rgb = AMBER; accent.line.fill.background(); accent.shadow.inherit = False
add_text(s, Inches(0.9), Inches(2.5), Inches(9), Inches(0.5), "DEMO PAYROLL SYSTEM", size=16, color=AMBER, bold=True)
add_text(s, Inches(0.9), Inches(2.95), Inches(11.5), Inches(1.6),
         "Payroll, Workforce & Accounting\nBuilt for Namibia's Security Industry",
         size=40, color=WHITE, bold=True, line_spacing=1.05)
add_text(s, Inches(0.9), Inches(4.55), Inches(10), Inches(0.6),
         "One platform for compliant payroll, site scheduling, client invoicing and live financials.",
         size=17, color=RGBColor(0xC7, 0xCE, 0xDA))
add_text(s, Inches(0.9), Inches(6.85), Inches(8), Inches(0.4),
         "Live demo tenant: Apex Shield Security", size=12, color=GRAY)

# ---------------------------------------------------------------- Slide 2: The problem
s = add_slide(); fill_bg(s, WHITE)
add_kicker(s, "The Problem")
add_text(s, Inches(0.6), Inches(0.75), Inches(10), Inches(0.8),
          "Security companies run payroll on spreadsheets—and it's costing them.",
          size=26, color=NAVY, bold=True)
add_bullets(s, Inches(0.6), Inches(1.9), Inches(11.8), Inches(4.8), [
    ("Compliance risk. ", "Minimum wage, SSC contributions, PAYE thresholds and the 70hr/week overtime cap (s.17(3)) are easy to get wrong manually."),
    ("No site-level visibility. ", "Guards rotate across dozens of client sites; spreadsheets can't track schedules, attendance and disciplinary history together."),
    ("Disconnected finances. ", "Payroll, client invoicing and the general ledger live in separate tools — nobody has a real-time view of cash, AR or margins."),
    ("Slow, error-prone payslips. ", "Manual tax and deduction calculations mean payday delays and disputes."),
], size=17, gap=18)
add_pagenum(s, 2)

# ---------------------------------------------------------------- Slide 3: Solution overview
s = add_slide(); fill_bg(s, WHITE)
add_kicker(s, "The Solution")
add_text(s, Inches(0.6), Inches(0.75), Inches(6.6), Inches(1.4),
          "One system: payroll, operations and accounting — fully connected.",
          size=24, color=NAVY, bold=True, line_spacing=1.05)
add_bullets(s, Inches(0.6), Inches(2.25), Inches(6.0), Inches(4.5), [
    "Live operational dashboard across employees, sites and pay periods",
    "Guided 7-step wizard: constants → sites → employees → schedule → attendance → payroll",
    "Real-time compliance reference built into every screen",
    "Role-based access for payroll, accounting and executive users",
], size=16, gap=14)
picture_framed(s, IMG("dashboard.png"), Inches(6.9), Inches(1.7), Inches(5.9), Inches(5.2))
add_pagenum(s, 3)

# ---------------------------------------------------------------- Slide 4: Workforce setup
s = add_slide(); fill_bg(s, WHITE)
add_kicker(s, "Workforce Operations")
add_text(s, Inches(0.6), Inches(0.75), Inches(11.5), Inches(0.7),
          "Employees and client sites, organized from day one.", size=24, color=NAVY, bold=True)
picture_framed(s, IMG("employees.png"), Inches(0.6), Inches(1.7), Inches(6.0), Inches(5.2))
picture_framed(s, IMG("sites.png"), Inches(6.8), Inches(1.7), Inches(6.0), Inches(5.2))
add_pagenum(s, 4)

# ---------------------------------------------------------------- Slide 5: Scheduling & attendance
s = add_slide(); fill_bg(s, WHITE)
add_kicker(s, "Workforce Operations")
add_text(s, Inches(0.6), Inches(0.75), Inches(11.5), Inches(0.7),
          "Shift scheduling and attendance, feeding straight into payroll.", size=24, color=NAVY, bold=True)
picture_framed(s, IMG("schedule.png"), Inches(0.6), Inches(1.7), Inches(6.0), Inches(5.2))
picture_framed(s, IMG("attendance.png"), Inches(6.8), Inches(1.7), Inches(6.0), Inches(5.2))
add_pagenum(s, 5)

# ---------------------------------------------------------------- Slide 6: Payroll engine
s = add_slide(); fill_bg(s, WHITE)
add_kicker(s, "Payroll Engine")
add_text(s, Inches(0.6), Inches(0.75), Inches(6.6), Inches(1.4),
          "Compliant payroll runs, calculated automatically.", size=24, color=NAVY, bold=True, line_spacing=1.05)
add_bullets(s, Inches(0.6), Inches(2.1), Inches(6.0), Inches(4.6), [
    ("PAYE & SSC handled. ", "Tax thresholds and social security contributions calculated per employee, every pay period."),
    ("Overtime cap enforced. ", "70hr/week cap (s.17(3)) tracked automatically against logged attendance."),
    ("One-click payslips. ", "Finalizing a pay period generates payslips and posts straight to the ledger."),
    ("Full audit trail. ", "Every pay run, adjustment and approval is logged for compliance review."),
], size=16, gap=16)
picture_framed(s, IMG("payroll.png"), Inches(6.9), Inches(1.7), Inches(5.9), Inches(5.2))
add_pagenum(s, 6)

# ---------------------------------------------------------------- Slide 7: Clients & invoicing
s = add_slide(); fill_bg(s, WHITE)
add_kicker(s, "Client Billing")
add_text(s, Inches(0.6), Inches(0.75), Inches(11.5), Inches(0.7),
          "Clients, sites and invoices — connected end to end.", size=24, color=NAVY, bold=True)
picture_framed(s, IMG("clients.png"), Inches(0.6), Inches(1.7), Inches(6.0), Inches(5.2))
picture_framed(s, IMG("invoices.png"), Inches(6.8), Inches(1.7), Inches(6.0), Inches(5.2))
add_pagenum(s, 7)

# ---------------------------------------------------------------- Slide 8: Accounting / financials
s = add_slide(); fill_bg(s, WHITE)
add_kicker(s, "Accounting")
add_text(s, Inches(0.6), Inches(0.75), Inches(6.6), Inches(1.4),
          "A live double-entry ledger — not a spreadsheet.", size=24, color=NAVY, bold=True, line_spacing=1.05)
add_bullets(s, Inches(0.6), Inches(2.1), Inches(6.0), Inches(4.6), [
    ("Auto-posted. ", "Invoices, payments and finalized payroll post to the ledger automatically — AR, AP, wages payable, tax control accounts."),
    ("Real-time P&L. ", "Revenue, expenses and net profit visible at a glance, always current."),
    ("Trial balance & AR aging. ", "Built-in reports for month-end close and collections follow-up."),
    ("AP & vendor bills. ", "Track payables alongside receivables in the same ledger."),
], size=16, gap=16)
picture_framed(s, IMG("accounting-pnl.png"), Inches(6.9), Inches(1.7), Inches(5.9), Inches(5.2))
add_pagenum(s, 8)

# ---------------------------------------------------------------- Slide 9: AI Executive Assistant
s = add_slide(); fill_bg(s, NAVY_LIGHT)
add_kicker(s, "AI Executive Assistant", color=AMBER)
add_text(s, Inches(0.6), Inches(0.75), Inches(6.2), Inches(1.5),
          "Ask your business a question — get a straight answer.", size=24, color=WHITE, bold=True, line_spacing=1.05)
add_bullets(s, Inches(0.6), Inches(2.2), Inches(5.9), Inches(4.4), [
    ("Plain-language queries. ", "“What's our margin on Apex contracts this month?” — answered from live financial and operational data."),
    ("Tenant-aware. ", "Every answer is scoped to your company's real numbers, not generic advice."),
    ("Executive-only. ", "Reserved for CEO/admin roles, alongside full clients and P&L visibility."),
], size=16, gap=16, color=RGBColor(0xE3,0xE7,0xEE))
picture_framed(s, IMG("ai-assistant.png"), Inches(6.9), Inches(1.7), Inches(5.9), Inches(5.2))
add_pagenum(s, 9)

# ---------------------------------------------------------------- Slide 10: Why it matters
s = add_slide(); fill_bg(s, WHITE)
add_kicker(s, "Why It Matters")
add_text(s, Inches(0.6), Inches(0.75), Inches(11), Inches(0.7),
          "Less risk, less admin, more visibility.", size=26, color=NAVY, bold=True)
cards = [
    ("Compliance built in", "Statutory rates and caps are encoded once, applied every pay run."),
    ("Hours saved monthly", "No more manual payslip math or reconciling three separate tools."),
    ("Real-time financial view", "Revenue, AR and margins visible the moment work is invoiced."),
]
cw = Inches(3.85); gap = Inches(0.3); left0 = Inches(0.6); top = Inches(2.1); ch = Inches(3.6)
for i, (title, desc) in enumerate(cards):
    l = left0 + i * (cw + gap)
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, top, cw, ch)
    card.fill.solid(); card.fill.fore_color.rgb = LIGHT_BG
    card.line.color.rgb = RGBColor(0xE2, 0xE6, 0xEB); card.line.width = Pt(1)
    card.shadow.inherit = False
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, top, cw, Inches(0.12))
    bar.fill.solid(); bar.fill.fore_color.rgb = AMBER; bar.line.fill.background(); bar.shadow.inherit = False
    add_text(s, l + Inches(0.3), top + Inches(0.45), cw - Inches(0.6), Inches(0.9), title,
              size=18, color=NAVY, bold=True, line_spacing=1.05)
    add_text(s, l + Inches(0.3), top + Inches(1.35), cw - Inches(0.6), Inches(2.0), desc,
              size=14, color=GRAY, line_spacing=1.2)
add_pagenum(s, 10)

# ---------------------------------------------------------------- Slide 11: Closing / CTA
s = add_slide(); fill_bg(s, NAVY)
accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(0.12))
accent.fill.solid(); accent.fill.fore_color.rgb = AMBER; accent.line.fill.background(); accent.shadow.inherit = False
add_text(s, Inches(0.9), Inches(2.6), Inches(11), Inches(1.2),
          "Let's set up your security company on this platform.", size=32, color=WHITE, bold=True, line_spacing=1.1)
add_text(s, Inches(0.9), Inches(3.9), Inches(10), Inches(0.6),
          "Payroll, scheduling, client invoicing and live financials — in one place, from day one.",
          size=16, color=RGBColor(0xC7, 0xCE, 0xDA))
add_text(s, Inches(0.9), Inches(6.6), Inches(8), Inches(0.4), "Demo Payroll System · Apex Shield Security demo tenant",
          size=11, color=GRAY)

prs.save(OUT)
print("Saved:", OUT)
